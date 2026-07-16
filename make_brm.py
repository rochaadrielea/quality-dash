"""
make_brm.py
-----------
Generates a 4-slide BRM PowerPoint from the SQLite database.

Layout mirrors Simone's BRM_Quality_BGLQC_May_26.pptx:
    Slide 1 - Cover
    Slide 2 - Quality Dashboard KPIs
    Slide 3 - Quality Performance & Trends (1/2)
    Slide 4 - Quality Performance & Trends (2/2)

Note: Numbers are auto-filled from SQLite. Manual review is still needed
before sending to Simone (colors, chart styling, comments sections).

Usage:
    from make_brm import build_brm_pptx
    pptx_bytes = build_brm_pptx(month="June 2026")
"""

from datetime import datetime
from io import BytesIO
import sqlite3

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION

DB_FILE = "quality.db"

# Beyond Gravity-ish colors (based on Simone's May deck)
BG_NAVY = RGBColor(0x1E, 0x27, 0x61)
BG_ORANGE = RGBColor(0xF2, 0x6E, 0x21)
BG_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
BG_GREEN = RGBColor(0x4C, 0xAF, 0x50)
BG_RED = RGBColor(0xE5, 0x3E, 0x3E)
BG_DARK = RGBColor(0x33, 0x33, 0x33)
BG_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def _run_query(sql):
    with sqlite3.connect(DB_FILE) as conn:
        return pd.read_sql(sql, conn)


def _add_text(slide, left, top, width, height, text, *,
              font_size=14, bold=False, color=BG_DARK, align_center=False, fill=None):
    """Helper to add a text box."""
    from pptx.enum.text import PP_ALIGN
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    tf = box.text_frame
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.03)
    tf.margin_bottom = Inches(0.03)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align_center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return box


def _add_kpi_card(slide, left, top, width, height, label, value, sub=None, color=BG_NAVY):
    """Big number card: colored strip + value + label."""
    # background card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = BG_WHITE
    card.line.color.rgb = color
    card.line.width = Pt(1.5)
    card.shadow.inherit = False

    # Label (top)
    _add_text(slide, left + Emu(50000), top + Emu(50000),
              width - Emu(100000), Inches(0.3),
              label, font_size=10, bold=True, color=color, align_center=True)

    # Big value (center)
    _add_text(slide, left, top + Inches(0.3),
              width, height - Inches(0.5),
              value, font_size=32, bold=True, color=BG_DARK, align_center=True)

    # Sub (bottom)
    if sub:
        _add_text(slide, left, top + height - Inches(0.3),
                  width, Inches(0.3),
                  sub, font_size=9, color=BG_DARK, align_center=True)


# ============================================================
# SLIDE BUILDERS
# ============================================================

def build_slide_1_cover(prs, month_label):
    """Cover slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Top navy band
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    band.fill.solid()
    band.fill.fore_color.rgb = BG_NAVY
    band.line.fill.background()

    _add_text(slide, Inches(0.5), Inches(0.3), Inches(12), Inches(0.6),
              "Beyond Gravity", font_size=24, bold=True, color=BG_WHITE)

    # Main title
    _add_text(slide, Inches(0.5), Inches(2.5), Inches(12), Inches(1),
              "BU Launchers Switzerland - Quality",
              font_size=44, bold=True, color=BG_NAVY)

    # Subtitle
    _add_text(slide, Inches(0.5), Inches(3.6), Inches(12), Inches(0.8),
              f"Business Review – {month_label}",
              font_size=28, color=BG_ORANGE)

    # Footer
    _add_text(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.3),
              f"Generated {datetime.now().strftime('%d.%m.%Y')}",
              font_size=10, color=BG_DARK)


def build_slide_2_dashboard(prs, month_label):
    """Quality dashboard with KPI cards."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title bar
    _add_text(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5),
              f"Quality Dashboard – {month_label}",
              font_size=24, bold=True, color=BG_NAVY)

    # Headline numbers
    kpis = _run_query("""
        SELECT
            SUM(CASE WHEN is_open=1 THEN 1 ELSE 0 END) AS ncs_wip,
            SUM(CASE WHEN is_open=0 THEN 1 ELSE 0 END) AS ncs_closed,
            SUM(CASE WHEN classification LIKE 'Major%' AND is_open=1 THEN 1 ELSE 0 END) AS major_open,
            SUM(CASE WHEN is_supplier_nc=1 AND is_open=1 THEN 1 ELSE 0 END) AS supplier_open,
            SUM(CASE WHEN is_supplier_nc=0 AND is_open=1 THEN 1 ELSE 0 END) AS production_open,
            SUM(CASE WHEN detection_area LIKE '%Customer%' AND is_open=1 THEN 1 ELSE 0 END) AS customer_complaints
        FROM nc
    """).iloc[0]

    # Section header - QUALITY
    _add_text(slide, Inches(0.4), Inches(0.9), Inches(4), Inches(0.3),
              "QUALITY", font_size=14, bold=True, color=BG_NAVY, fill=BG_LIGHT)

    # 5 KPI cards in a row
    card_w = Inches(2.5)
    card_h = Inches(1.2)
    top_row = Inches(1.3)
    positions = [Inches(0.4 + i * 2.6) for i in range(5)]

    _add_kpi_card(slide, positions[0], top_row, card_w, card_h,
                  "NCs OPENED (WIP)", int(kpis["ncs_wip"] or 0), sub="current month")
    _add_kpi_card(slide, positions[1], top_row, card_w, card_h,
                  "NCs CLOSED", int(kpis["ncs_closed"] or 0), sub="current month", color=BG_GREEN)
    _add_kpi_card(slide, positions[2], top_row, card_w, card_h,
                  "MAJOR NCs (OPEN)", int(kpis["major_open"] or 0), color=BG_RED)
    _add_kpi_card(slide, positions[3], top_row, card_w, card_h,
                  "PRODUCTION NCs", int(kpis["production_open"] or 0))
    _add_kpi_card(slide, positions[4], top_row, card_w, card_h,
                  "SUPPLIER NCs", int(kpis["supplier_open"] or 0), color=BG_ORANGE)

    # Second section - DELIVERY / COST / PEOPLE placeholders (Simone fills these)
    _add_text(slide, Inches(0.4), Inches(2.7), Inches(4), Inches(0.3),
              "DELIVERY", font_size=14, bold=True, color=BG_NAVY, fill=BG_LIGHT)

    _add_text(slide, Inches(0.4), Inches(3.05), Inches(12.5), Inches(0.7),
              "Flow Attainment: [manual] · CAPA Status WIP/Delayed/Closed: [manual] · "
              "MRR/PRR: [manual] · DRBs: [manual]",
              font_size=11, color=BG_DARK)

    _add_text(slide, Inches(0.4), Inches(3.9), Inches(4), Inches(0.3),
              "COST", font_size=14, bold=True, color=BG_NAVY, fill=BG_LIGHT)
    _add_text(slide, Inches(0.4), Inches(4.2), Inches(12.5), Inches(0.4),
              "CoPQ Monthly: [manual, target 3%] · CoPQ YTD: [manual, target 3%]",
              font_size=11, color=BG_DARK)

    _add_text(slide, Inches(0.4), Inches(4.7), Inches(4), Inches(0.3),
              "PEOPLE", font_size=14, bold=True, color=BG_NAVY, fill=BG_LIGHT)
    _add_text(slide, Inches(0.4), Inches(5.0), Inches(12.5), Inches(0.4),
              "FTE: [manual, target 2.5] · Direct Rate: [manual]",
              font_size=11, color=BG_DARK)

    # Highlights / Lowlights / Outlook boxes
    box_top = Inches(5.6)
    box_h = Inches(1.6)
    box_w = Inches(3.0)
    labels_boxes = [
        ("HIGHLIGHTS", Inches(0.4), BG_GREEN),
        ("LOWLIGHTS", Inches(3.6), BG_RED),
        ("OUTLOOK", Inches(6.8), BG_NAVY),
        ("KEY ISSUES / RISKS", Inches(10.0), BG_ORANGE),
    ]
    for label, left, color in labels_boxes:
        card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, box_top, box_w, box_h)
        card.fill.solid()
        card.fill.fore_color.rgb = BG_WHITE
        card.line.color.rgb = color
        card.line.width = Pt(1)
        _add_text(slide, left, box_top, box_w, Inches(0.3),
                  label, font_size=11, bold=True, color=color, align_center=True)
        _add_text(slide, left + Inches(0.1), box_top + Inches(0.35),
                  box_w - Inches(0.2), box_h - Inches(0.4),
                  "[Manual comments to be added]",
                  font_size=10, color=BG_DARK)


def build_slide_3_trends1(prs, month_label):
    """Slide 3: NCs by Project + NCs by Area."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_text(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5),
              f"Quality Performance & Trends (1/2) – {month_label}",
              font_size=22, bold=True, color=BG_NAVY)

    # Top 6 Projects WIP
    top6_wip = _run_query("""
        SELECT COALESCE(project, '(no project)') AS project, COUNT(*) AS n
        FROM nc WHERE is_open = 1
        GROUP BY project ORDER BY n DESC LIMIT 6
    """)

    _add_text(slide, Inches(0.4), Inches(0.9), Inches(6), Inches(0.3),
              "Top 6 Projects – NCs WIP", font_size=13, bold=True, color=BG_NAVY)

    chart_data = CategoryChartData()
    chart_data.categories = top6_wip["project"].tolist()
    chart_data.add_series("Open NCs", top6_wip["n"].tolist())
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.4), Inches(1.3), Inches(6), Inches(3),
        chart_data,
    ).chart
    chart.has_title = False
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.font.size = Pt(9)

    # By Area (YTD)
    by_area = _run_query("""
        SELECT COALESCE(detection_area, 'BLANK - to clean') AS area, COUNT(*) AS n
        FROM nc WHERE substr(created_on, 1, 4) = '2026'
        GROUP BY area ORDER BY n DESC LIMIT 10
    """)

    _add_text(slide, Inches(7.0), Inches(0.9), Inches(6), Inches(0.3),
              "NCs by Detection Area (YTD 2026)",
              font_size=13, bold=True, color=BG_NAVY)

    chart_data2 = CategoryChartData()
    chart_data2.categories = by_area["area"].tolist()
    chart_data2.add_series("NCs", by_area["n"].tolist())
    chart2 = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(7.0), Inches(1.3), Inches(6), Inches(3),
        chart_data2,
    ).chart
    chart2.has_title = False
    chart2.has_legend = False
    plot2 = chart2.plots[0]
    plot2.has_data_labels = True
    plot2.data_labels.font.size = Pt(9)

    # Top 6 Projects YTD
    top6_ytd = _run_query("""
        SELECT COALESCE(project, '(no project)') AS project, COUNT(*) AS n
        FROM nc WHERE substr(created_on, 1, 4) = '2026'
        GROUP BY project ORDER BY n DESC LIMIT 6
    """)

    _add_text(slide, Inches(0.4), Inches(4.4), Inches(6), Inches(0.3),
              "Top 6 Projects – NCs Opened YTD 2026",
              font_size=13, bold=True, color=BG_NAVY)

    chart_data3 = CategoryChartData()
    chart_data3.categories = top6_ytd["project"].tolist()
    chart_data3.add_series("Opened YTD", top6_ytd["n"].tolist())
    chart3 = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.4), Inches(4.8), Inches(6), Inches(2.4),
        chart_data3,
    ).chart
    chart3.has_title = False
    chart3.has_legend = False
    chart3.plots[0].has_data_labels = True
    chart3.plots[0].data_labels.font.size = Pt(9)

    # NC total per year
    per_year = _run_query("""
        SELECT substr(created_on, 1, 4) AS yr, COUNT(*) AS n
        FROM nc WHERE created_on IS NOT NULL
        GROUP BY yr ORDER BY yr
    """)

    _add_text(slide, Inches(7.0), Inches(4.4), Inches(6), Inches(0.3),
              "NC Total per Year", font_size=13, bold=True, color=BG_NAVY)

    chart_data4 = CategoryChartData()
    chart_data4.categories = per_year["yr"].tolist()
    chart_data4.add_series("NCs", per_year["n"].tolist())
    chart4 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.0), Inches(4.8), Inches(6), Inches(2.4),
        chart_data4,
    ).chart
    chart4.has_title = False
    chart4.has_legend = False
    chart4.plots[0].has_data_labels = True
    chart4.plots[0].data_labels.font.size = Pt(9)


def build_slide_4_trends2(prs, month_label):
    """Slide 4: Production vs Supplier + monthly trend."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _add_text(slide, Inches(0.4), Inches(0.2), Inches(12.5), Inches(0.5),
              f"Quality Performance & Trends (2/2) – {month_label}",
              font_size=22, bold=True, color=BG_NAVY)

    # Monthly trend (opens vs closes)
    monthly = _run_query("""
        WITH opens AS (
            SELECT substr(created_on, 1, 7) AS month, COUNT(*) AS opened
            FROM nc WHERE created_on IS NOT NULL GROUP BY month
        ),
        closes AS (
            SELECT substr(closure_date, 1, 7) AS month, COUNT(*) AS closed
            FROM nc WHERE closure_date IS NOT NULL GROUP BY month
        )
        SELECT o.month, o.opened, COALESCE(c.closed, 0) AS closed
        FROM opens o LEFT JOIN closes c USING (month)
        WHERE o.month >= '2026-01'
        ORDER BY o.month
    """)

    _add_text(slide, Inches(0.4), Inches(0.9), Inches(12), Inches(0.3),
              "Monthly Trend – Opens vs Closes (2026)",
              font_size=13, bold=True, color=BG_NAVY)

    chart_data = CategoryChartData()
    chart_data.categories = monthly["month"].tolist()
    chart_data.add_series("Opened", monthly["opened"].tolist())
    chart_data.add_series("Closed", monthly["closed"].tolist())
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.4), Inches(1.3), Inches(12.5), Inches(2.5),
        chart_data,
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False

    # Production vs Supplier
    prod_sup = _run_query("""
        SELECT
            CASE WHEN is_supplier_nc = 1 THEN 'Supplier' ELSE 'Production' END AS source,
            SUM(CASE WHEN is_open = 1 THEN 1 ELSE 0 END) AS open_ncs,
            SUM(CASE WHEN is_open = 0 THEN 1 ELSE 0 END) AS closed_ncs
        FROM nc GROUP BY source
    """)

    _add_text(slide, Inches(0.4), Inches(4.0), Inches(6), Inches(0.3),
              "Production vs Supplier NCs",
              font_size=13, bold=True, color=BG_NAVY)

    chart_data2 = CategoryChartData()
    chart_data2.categories = prod_sup["source"].tolist()
    chart_data2.add_series("Open", prod_sup["open_ncs"].tolist())
    chart_data2.add_series("Closed", prod_sup["closed_ncs"].tolist())
    chart2 = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.4), Inches(4.4), Inches(6), Inches(2.8),
        chart_data2,
    ).chart
    chart2.has_legend = True
    chart2.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart2.plots[0].has_data_labels = True
    chart2.plots[0].data_labels.font.size = Pt(10)

    # Owner workload table
    owners = _run_query("""
        SELECT COALESCE(owner, '(no owner)') AS owner,
               SUM(CASE WHEN is_open = 1 THEN 1 ELSE 0 END) AS open_ncs,
               MAX(days_open) AS oldest_days
        FROM nc
        GROUP BY owner
        HAVING open_ncs > 0
        ORDER BY open_ncs DESC
        LIMIT 8
    """)

    _add_text(slide, Inches(7.0), Inches(4.0), Inches(6), Inches(0.3),
              "Open NCs by Owner (Top 8)",
              font_size=13, bold=True, color=BG_NAVY)

    # Simple table
    rows = len(owners) + 1
    cols = 3
    table_shape = slide.shapes.add_table(rows, cols,
                                          Inches(7.0), Inches(4.4),
                                          Inches(6), Inches(2.8))
    table = table_shape.table
    hdrs = ["Owner", "Open NCs", "Oldest (days)"]
    for i, h in enumerate(hdrs):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_NAVY
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = BG_WHITE
                r.font.bold = True
                r.font.size = Pt(11)

    for i, row in enumerate(owners.itertuples(), start=1):
        table.cell(i, 0).text = str(row.owner)
        table.cell(i, 1).text = str(int(row.open_ncs))
        table.cell(i, 2).text = str(int(row.oldest_days)) if pd.notna(row.oldest_days) else ""
        for c in range(cols):
            for p in table.cell(i, c).text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(10)


# ============================================================
# MAIN ENTRY
# ============================================================

def build_brm_pptx(month="June 2026") -> bytes:
    """Build the full 4-slide BRM. Returns pptx as bytes."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    build_slide_1_cover(prs, month)
    build_slide_2_dashboard(prs, month)
    build_slide_3_trends1(prs, month)
    build_slide_4_trends2(prs, month)

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


if __name__ == "__main__":
    data = build_brm_pptx("June 2026")
    with open("BRM_Quality_June_2026.pptx", "wb") as f:
        f.write(data)
    print(f"✓ Written BRM_Quality_June_2026.pptx ({len(data)/1024:.1f} KB)")
